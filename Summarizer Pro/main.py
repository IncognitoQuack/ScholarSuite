import sys
import os
import re
import textwrap
from PyQt5.QtWidgets import (
    QApplication, QMainWindow, QAction, QFileDialog,
    QHBoxLayout, QVBoxLayout, QWidget, QFrame, QLabel,
    QLineEdit, QPushButton, QTextEdit, QComboBox, QMessageBox, QCheckBox
)
from PyQt5.QtCore import Qt
from PyQt5.QtGui import QFont

import PyPDF2
import docx2txt

# PPTX Imports
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# ReportLab Imports
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, PageBreak, Image, ListFlowable, ListItem
)
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.pagesizes import LETTER
from reportlab.lib.enums import TA_LEFT
from reportlab.lib.colors import black as BLACK

########################################
# PPT Styling Constants
########################################

PRIMARY_COLOR = RGBColor(0, 70, 122)   # Dark blue
TITLE_FONT = "Calibri"
BODY_FONT = "Calibri"

########################################
# PPT Helper Functions
########################################

def add_footer(slide, footer_text, prs):
    """Adds a simple footer textbox at the bottom-right corner of the slide."""
    width = Inches(3)
    height = Inches(0.3)
    left = prs.slide_width - width - Inches(0.3)
    top = prs.slide_height - height - Inches(0.1)

    footer_box = slide.shapes.add_textbox(left, top, width, height)
    footer_tf = footer_box.text_frame
    footer_p = footer_tf.paragraphs[0]
    footer_p.text = footer_text
    footer_p.font.name = BODY_FONT
    footer_p.font.size = Pt(12)
    footer_p.font.color.rgb = RGBColor(80, 80, 80)
    footer_p.alignment = PP_ALIGN.RIGHT

def add_title_shape(slide, title_text, width, height, left, top, fill_color):
    """Adds a colored rectangular shape with centered title text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color  # Hide border

    text_frame = shape.text_frame
    text_frame.text = title_text
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.bold = True
        paragraph.font.size = Pt(32)
        paragraph.font.name = TITLE_FONT
        paragraph.font.color.rgb = RGBColor(255, 255, 255)

def add_subheading(slide, subheading_text, left, top, width, height):
    """Adds a subheading text box."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = subheading_text
    p.font.name = BODY_FONT
    p.font.size = Pt(18)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    p.font.color.rgb = RGBColor(0, 0, 0)

def add_bullet_points(slide, bullet_list, left, top, width, height):
    """
    Adds bullet points to a text box, wrapping lines so they don't overflow.
    We'll use a fixed wrap width (e.g. ~100 chars). You can adjust as needed.
    """
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.clear()

    WRAP_WIDTH = 80  # character-based wrapping

    for item in bullet_list:
        # Wrap each bullet line if it's too long
        wrapped_lines = textwrap.wrap(item, WRAP_WIDTH)
        if not wrapped_lines:
            wrapped_lines = [""]

        # The first wrapped line is the bullet; subsequent lines become sub-bullet lines
        first_line = True
        for wline in wrapped_lines:
            p = tf.add_paragraph()
            p.font.size = Pt(18)
            p.font.name = BODY_FONT
            p.font.color.rgb = RGBColor(0, 0, 0)
            if first_line:
                p.level = 0
                p.text = wline
                first_line = False
            else:
                p.level = 1
                p.text = wline

def add_image(slide, image_path, left, top, width, height):
    """Adds an image to the slide if desired."""
    slide.shapes.add_picture(image_path, left, top, width, height)

def create_title_slide(prs, main_title, subtitle, footer_text):
    """Creates a stylized title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title shape across the top
    add_title_shape(
        slide=slide,
        title_text=main_title,
        width=prs.slide_width,
        height=Inches(2),
        left=Inches(0),
        top=Inches(0),
        fill_color=PRIMARY_COLOR
    )

    # Subtitle below the shape
    sub_left = Inches(1)
    sub_top = Inches(2.3)
    sub_width = prs.slide_width - Inches(2)
    sub_height = Inches(1.5)
    if subtitle:
        add_subheading(slide, subtitle, sub_left, sub_top, sub_width, sub_height)

    add_footer(slide, footer_text, prs)

def create_stylized_slide(prs, slide_title, subheading, bullet_points, footer_text, image_path=None):
    """
    Creates a stylized slide with:
      1) A colored title shape on the left.
      2) A subheading at the top-right.
      3) Bullet points in the lower-right area.
      4) An optional image on the right or left.
    """
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # 1) Title shape on the left
    shape_width = Inches(3)
    shape_height = prs.slide_height
    add_title_shape(
        slide=slide,
        title_text=slide_title,
        width=shape_width,
        height=shape_height,
        left=Inches(0),
        top=Inches(0),
        fill_color=PRIMARY_COLOR
    )

    # 2) Subheading at the top-right
    sub_left = shape_width + Inches(0.3)
    sub_top = Inches(0.2)
    sub_width = prs.slide_width - shape_width - Inches(0.6)
    sub_height = Inches(1)
    if subheading:
        add_subheading(slide, subheading, sub_left, sub_top, sub_width, sub_height)

    # 3) Bullet points in the lower-right area
    bullet_left = shape_width + Inches(0.3)
    bullet_top = Inches(1.2)
    bullet_width = prs.slide_width - shape_width - Inches(0.6)
    bullet_height = prs.slide_height - Inches(1.8)
    if bullet_points:
        add_bullet_points(slide, bullet_points, bullet_left, bullet_top, bullet_width, bullet_height)

    # 4) Optional image
    if image_path:
        img_left = prs.slide_width - Inches(3.5)
        img_top = Inches(1)
        img_width = Inches(3)
        img_height = Inches(2)
        add_image(slide, image_path, img_left, img_top, img_width, img_height)

    add_footer(slide, footer_text, prs)

########################################
# The main PyQt application
########################################

class MainWindow(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("StylishSummarizer Pro")
        self.resize(1300, 800)
        self.outputPath = ""
        self.image_files = {}  # For <--- img n ---> placeholders
        self.setupUI()

    def setupUI(self):
        # Menu bar
        menubar = self.menuBar()
        helpMenu = menubar.addMenu("Help")
        helpAction = QAction("Usage Instructions", self)
        helpAction.triggered.connect(self.show_help)
        helpMenu.addAction(helpAction)

        # Main container
        centralWidget = QWidget()
        mainLayout = QHBoxLayout(centralWidget)
        self.setCentralWidget(centralWidget)

        # Left panel: Prompt Builder
        leftPanel = QFrame()
        leftLayout = QVBoxLayout(leftPanel)
        leftPanel.setObjectName("leftPanel")

        titleLabel = QLabel("Prompt Builder")
        titleLabel.setFont(QFont("Arial", 12, QFont.Bold))
        leftLayout.addWidget(titleLabel)

        # Document Type
        self.docTypeCombo = QComboBox()
        self.docTypeCombo.addItems(["Case Study", "Report", "Project Info", "Notes"])
        docTypeLayout = QHBoxLayout()
        docTypeLayout.addWidget(QLabel("Document Type:"))
        docTypeLayout.addWidget(self.docTypeCombo)
        leftLayout.addLayout(docTypeLayout)

        # Length
        self.lengthCombo = QComboBox()
        self.lengthCombo.addItems(["Short", "Mid", "Long", "Very-Long"])
        lengthLayout = QHBoxLayout()
        lengthLayout.addWidget(QLabel("Length:"))
        lengthLayout.addWidget(self.lengthCombo)
        leftLayout.addLayout(lengthLayout)

        # Info Type
        self.infoTypeCombo = QComboBox()
        self.infoTypeCombo.addItems(["I have a topic name only", "I have a file to attach"])
        self.infoTypeCombo.currentIndexChanged.connect(self.toggleTopicInput)
        infoTypeLayout = QHBoxLayout()
        infoTypeLayout.addWidget(QLabel("Info Type:"))
        infoTypeLayout.addWidget(self.infoTypeCombo)
        leftLayout.addLayout(infoTypeLayout)

        # Topic name
        self.topicInput = QLineEdit()
        self.topicInput.setPlaceholderText("Enter the topic name here...")
        leftLayout.addWidget(self.topicInput)

        # Additional instructions
        self.additionalInstructions = QLineEdit()
        self.additionalInstructions.setPlaceholderText("Any extra instructions for the LLM prompt...")
        leftLayout.addWidget(self.additionalInstructions)

        # Copy prompt button
        self.copyPromptButton = QPushButton("Copy Prompt to Clipboard")
        self.copyPromptButton.clicked.connect(self.copy_prompt_to_clipboard)
        leftLayout.addWidget(self.copyPromptButton)

        leftLayout.addStretch()

        # Middle panel: Text Editor
        middlePanel = QFrame()
        middleLayout = QVBoxLayout(middlePanel)
        middlePanel.setObjectName("middlePanel")

        editorLabel = QLabel("Document Editor")
        editorLabel.setFont(QFont("Arial", 12, QFont.Bold))
        middleLayout.addWidget(editorLabel)

        self.textEdit = QTextEdit()
        self.textEdit.setAcceptDrops(True)
        self.textEdit.setPlaceholderText("Paste text or drag & drop PDF, DOCX, or TXT files here...")
        middleLayout.addWidget(self.textEdit)

        # Right panel: Output Options
        rightPanel = QFrame()
        rightLayout = QVBoxLayout(rightPanel)
        rightPanel.setObjectName("rightPanel")

        optionsLabel = QLabel("Output Options")
        optionsLabel.setFont(QFont("Arial", 12, QFont.Bold))
        rightLayout.addWidget(optionsLabel)

        # File Name
        fileNameLayout = QHBoxLayout()
        fileNameLayout.addWidget(QLabel("File Name:"))
        self.fileNameInput = QLineEdit()
        self.fileNameInput.setPlaceholderText("Output file name (no extension)")
        fileNameLayout.addWidget(self.fileNameInput)
        rightLayout.addLayout(fileNameLayout)

        # Output directory
        self.selectDirButton = QPushButton("Select Output Directory")
        self.selectDirButton.clicked.connect(self.selectOutputDirectory)
        rightLayout.addWidget(self.selectDirButton)

        # Format
        formatLayout = QHBoxLayout()
        formatLayout.addWidget(QLabel("Format:"))
        self.formatCombo = QComboBox()
        self.formatCombo.addItems(["PPT", "PDF"])
        formatLayout.addWidget(self.formatCombo)
        rightLayout.addLayout(formatLayout)

        # Custom Watermark
        self.customWatermarkCheck = QCheckBox("Use Custom Watermark?")
        self.customWatermarkCheck.setChecked(False)
        self.customWatermarkCheck.stateChanged.connect(self.toggleWatermarkInput)
        rightLayout.addWidget(self.customWatermarkCheck)

        self.customWatermarkInput = QLineEdit()
        self.customWatermarkInput.setPlaceholderText("Enter your name or group name")
        self.customWatermarkInput.setEnabled(False)
        rightLayout.addWidget(self.customWatermarkInput)

        # Generate button
        self.generateButton = QPushButton("Generate")
        self.generateButton.clicked.connect(self.generateOutput)
        rightLayout.addWidget(self.generateButton)

        mainLayout.addWidget(leftPanel, stretch=2)
        mainLayout.addWidget(middlePanel, stretch=3)
        mainLayout.addWidget(rightPanel, stretch=2)

        self.setDarkTheme()

    def setDarkTheme(self):
        """Apply a dark QSS to give a professional, modern look."""
        dark_stylesheet = """
        QMainWindow {
            background-color: #2b2b2b;
        }
        QMenuBar {
            background-color: #2b2b2b;
            color: #ffffff;
        }
        QMenuBar::item {
            background-color: #2b2b2b;
            color: #ffffff;
        }
        QMenuBar::item:selected {
            background-color: #3c3c3c;
        }
        QMenu {
            background-color: #2b2b2b;
            color: #ffffff;
        }
        QMenu::item:selected {
            background-color: #3c3c3c;
        }
        QLabel {
            color: #ffffff;
        }
        QLineEdit, QTextEdit, QComboBox {
            background-color: #3c3c3c;
            color: #ffffff;
            border: 1px solid #5a5a5a;
            border-radius: 4px;
            padding: 4px;
        }
        QFrame#leftPanel, QFrame#middlePanel, QFrame#rightPanel {
            background-color: #2f2f2f;
            border: 1px solid #444;
            border-radius: 6px;
        }
        QPushButton {
            background-color: #3c3c3c;
            color: #ffffff;
            border: 1px solid #5a5a5a;
            border-radius: 4px;
            padding: 6px;
        }
        QPushButton:hover {
            background-color: #505050;
        }
        QCheckBox {
            color: #ffffff;
        }
        """
        self.setStyleSheet(dark_stylesheet)

    def show_help(self):
        msg = QMessageBox(self)
        msg.setWindowTitle("Usage Instructions")
        msg.setText(
            "1. Use the Prompt Builder (left panel) to specify the document type, length, topic or file, etc.\n"
            "   Click 'Copy Prompt to Clipboard' to get a well-structured prompt for ChatGPT.\n\n"
            "2. In the Document Editor (middle panel), paste text or drag & drop PDF/DOCX/TXT files.\n"
            "   - Headings: *h1- Title* or *h2- Title*\n"
            "   - Bullet lines: '* bullet text'\n"
            "   - *bold*, _italic_, <--- img n ---> for images.\n\n"
            "3. In Output Options (right panel), set file name, directory, format (PPT or PDF), etc.\n"
            "   Check 'Use Custom Watermark?' if you want your own text instead of 'Generated by StylishSummarizer Pro'.\n"
            "4. Click 'Generate' to produce a stylized PPT (with auto line wrapping) or a structured PDF."
        )
        msg.exec_()

    def toggleTopicInput(self):
        if self.infoTypeCombo.currentText() == "I have a topic name only":
            self.topicInput.setVisible(True)
        else:
            self.topicInput.setVisible(False)

    def toggleWatermarkInput(self):
        if self.customWatermarkCheck.isChecked():
            self.customWatermarkInput.setEnabled(True)
        else:
            self.customWatermarkInput.setEnabled(False)

    def copy_prompt_to_clipboard(self):
        doc_type = self.docTypeCombo.currentText()
        length = self.lengthCombo.currentText()
        info_type = self.infoTypeCombo.currentText()
        topic = self.topicInput.text().strip()
        additional = self.additionalInstructions.text().strip()

        prompt = f"Write a detailed {doc_type} (length of content: {length.lower()}) "
        if info_type == "I have a topic name only":
            if topic:
                prompt += f'on the topic "{topic}" by researching relevant information from the web.'
            else:
                prompt += 'on the topic "SAMPLE_TOPIC" by researching relevant information from the web.'
        else:
            prompt += "for the topic I have uploaded/linked a file in this chat."

        prompt += (
            "\nThe content should be comprehensive and well-structured, following a formal case study format.\n"
            f"Please ensure the content is {length.lower()} and thorough, with the following guidelines for styling:\n\n"
            "Use h1 for main headings (Title).\n"
            "Use h2 for subheadings (Title).\n"
            "Bullet points should be formatted with * text* (a line starting with an asterisk and a space).\n"
            "Use bold for bold text and italic for italics.\n"
            "Where appropriate, include images with the format: <--- img n --->.\n"
            "The document should be written in a clear, detailed manner, ensuring a professional and structured approach.\n"
            "NOTE: don't use any code snippets in the answer; use normal text, but apply the guidelines above."
        )

        if additional:
            prompt += f"\nAdditional instructions:\n{additional}\n"

        QApplication.clipboard().setText(prompt)
        QMessageBox.information(self, "Prompt Copied", "Prompt has been copied to the clipboard.")

    def dragEnterEvent(self, event):
        if event.mimeData().hasUrls():
            event.acceptProposedAction()

    def dropEvent(self, event):
        for url in event.mimeData().urls():
            file_path = url.toLocalFile()
            text = ""
            if file_path.lower().endswith('.pdf'):
                text = self.extract_text_from_pdf(file_path)
            elif file_path.lower().endswith('.docx'):
                text = self.extract_text_from_docx(file_path)
            elif file_path.lower().endswith('.txt'):
                try:
                    with open(file_path, 'r', encoding='utf-8') as f:
                        text = f.read()
                except Exception as e:
                    print(f"Error reading TXT: {e}")
            if text:
                self.textEdit.append(text)

    def extract_text_from_pdf(self, file_path):
        text = ""
        try:
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
        except Exception as e:
            print(f"Error reading PDF: {e}")
        return text

    def extract_text_from_docx(self, file_path):
        try:
            return docx2txt.process(file_path)
        except Exception as e:
            print(f"Error reading DOCX: {e}")
            return ""

    def selectOutputDirectory(self):
        directory = QFileDialog.getExistingDirectory(self, "Select Output Directory")
        if directory:
            self.outputPath = directory

    def prompt_for_images(self, placeholders):
        """For each detected placeholder, prompt the user to select an image file."""
        for ph in placeholders:
            if ph not in self.image_files:
                file_path, _ = QFileDialog.getOpenFileName(
                    self,
                    f"Select image for placeholder {ph}",
                    "",
                    "Images (*.png *.jpg *.jpeg *.bmp)"
                )
                if file_path:
                    self.image_files[ph] = file_path

    def generateOutput(self):
        raw_text = self.textEdit.toPlainText()
        if not raw_text.strip():
            QMessageBox.warning(self, "Input Error", "Please provide some text or files first.")
            return

        file_name = self.fileNameInput.text().strip() or "output"
        output_format = self.formatCombo.currentText()
        doc_type = self.docTypeCombo.currentText()

        # If custom watermark is checked, we use that; else default
        if self.customWatermarkCheck.isChecked():
            custom_text = self.customWatermarkInput.text().strip()
            watermark_text = f"{doc_type} - {custom_text}" if custom_text else f"{doc_type} - Generated by StylishSummarizer Pro"
        else:
            watermark_text = f"{doc_type} - Generated by StylishSummarizer Pro"

        # Detect image placeholders
        placeholders = re.findall(r'<---\s*img\s+(\d+)\s*--->', raw_text)
        if placeholders:
            self.prompt_for_images(placeholders)

        try:
            if output_format == "PPT":
                out_file = self.generate_ppt(doc_type, file_name, raw_text, watermark_text)
            else:
                out_file = self.generate_pdf(doc_type, file_name, raw_text, watermark_text)

            QMessageBox.information(self, "Success", f"File generated successfully at:\n{out_file}")
        except Exception as e:
            QMessageBox.critical(self, "Generation Error", f"An error occurred: {e}")

    #####################################################
    # PPT Generation with auto line-wrapping & subheading
    #####################################################
    def generate_ppt(self, doc_type, file_name, raw_text, watermark_text):
        """
        - Title slide with doc_type + watermark
        - Parse text lines for *h1- (new slide), *h2- (subheading), bullet lines, etc.
        - Auto line wrap bullet lines, chunk slides if too many lines
        - Insert images as separate slides
        """
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # 1) Title slide
        create_title_slide(
            prs,
            main_title=doc_type,
            subtitle=watermark_text,
            footer_text=watermark_text
        )

        # 2) Parse lines => multiple slides
        lines = raw_text.split("\n")
        slides_data = []
        current_title = "Content"
        current_subheading = ""
        current_bullets = []

        MAX_BULLETS_PER_SLIDE = 6

        def push_slide():
            if current_title or current_subheading or current_bullets:
                slides_data.append({
                    "title": current_title or "Content",
                    "subheading": current_subheading,
                    "bullets": current_bullets[:]
                })

        for line in lines:
            line_stripped = line.strip()
            if not line_stripped:
                continue

            # Check for *h1-
            h1_match = re.match(r'^\*h1-\s+(.*)$', line_stripped, flags=re.IGNORECASE)
            if h1_match:
                # push old slide
                if current_title or current_subheading or current_bullets:
                    push_slide()
                    current_bullets.clear()
                current_title = h1_match.group(1)
                current_subheading = ""
                continue

            # Check for *h2-
            h2_match = re.match(r'^\*h2-\s+(.*)$', line_stripped, flags=re.IGNORECASE)
            if h2_match:
                # push old slide
                if current_title or current_subheading or current_bullets:
                    push_slide()
                    current_bullets.clear()
                current_subheading = h2_match.group(1)
                current_title = ""  # if there's no h1, we'll just do subheading
                continue

            # bullet line => starts with "* "
            bullet_match = re.match(r'^\*\s+(.*)$', line_stripped)
            if bullet_match:
                current_bullets.append(bullet_match.group(1))
            else:
                # normal line => treat as bullet
                current_bullets.append(line_stripped)

            # chunk slides if too many bullets
            if len(current_bullets) >= MAX_BULLETS_PER_SLIDE:
                push_slide()
                current_bullets.clear()
                current_title = "Content"
                current_subheading = ""

        # push the last slide if there's leftover
        if current_title or current_subheading or current_bullets:
            push_slide()

        # create stylized slides
        for sd in slides_data:
            create_stylized_slide(
                prs=prs,
                slide_title=sd["title"],
                subheading=sd["subheading"],
                bullet_points=sd["bullets"],
                footer_text=watermark_text
            )

        # 3) Insert images => separate slides
        for ph, img_path in self.image_files.items():
            create_stylized_slide(
                prs,
                slide_title="Image Slide",
                subheading=f"Placeholder {ph}",
                bullet_points=[],
                footer_text=watermark_text,
                image_path=img_path
            )

        # Save
        if not self.outputPath:
            out_file = file_name + ".pptx"
        else:
            out_file = os.path.join(self.outputPath, file_name + ".pptx")
        prs.save(out_file)
        return out_file

    #####################################################
    # PDF Generation with structured layout
    #####################################################
    def generate_pdf(self, doc_type, file_name, raw_text, watermark_text):
        """
        We parse headings, bullets, and normal lines for a professional layout:
          - h1 => big heading
          - h2 => smaller heading
          - bullet lines => bullet points
          - normal lines => paragraphs
        We also handle *bold* => <b> and _italic_ => <i>.
        Each image => new page.
        """
        if not self.outputPath:
            out_file = file_name + ".pdf"
        else:
            out_file = os.path.join(self.outputPath, file_name + ".pdf")

        doc = SimpleDocTemplate(
            out_file,
            pagesize=LETTER,
            rightMargin=72,
            leftMargin=72,
            topMargin=72,
            bottomMargin=72
        )

        styles = getSampleStyleSheet()

        styleH1 = ParagraphStyle('Heading1', parent=styles['Heading1'], fontSize=16, spaceAfter=10)
        styleH2 = ParagraphStyle('Heading2', parent=styles['Heading2'], fontSize=14, spaceAfter=8)
        styleNormal = ParagraphStyle('Normal', parent=styles['Normal'], fontSize=11, leading=14)
        styleNormal.allowHTML = True

        flowables = []

        # Title page
        flowables.append(Paragraph(f"{doc_type}", styleH1))
        flowables.append(Paragraph(f"{watermark_text}", styleNormal))
        flowables.append(Spacer(1, 40))
        flowables.append(PageBreak())

        # Minimal transforms for bold & italic
        def transform_bold_italic(line):
            # *bold* => <b>bold</b>
            line = re.sub(r'\*(?!\s)(.*?)\*', r'<b>\1</b>', line)
            # _italic_ => <i>italic</i>
            line = re.sub(r'_(?!\s)(.*?)_', r'<i>\1</i>', line)
            return line

        # We'll parse line by line: h1 => big heading, h2 => smaller heading, bullet => bullet item, else paragraph
        lines = raw_text.splitlines()
        bullet_buffer = []  # store bullet lines, flush them as a ListFlowable
        def flush_bullets():
            if bullet_buffer:
                # create a bullet list
                items = []
                for bline in bullet_buffer:
                    items.append(ListItem(Paragraph(transform_bold_italic(bline), styleNormal), bulletColor=BLACK))
                flowables.append(ListFlowable(items, bulletType='bullet'))
                flowables.append(Spacer(1, 10))
                bullet_buffer.clear()

        for line in lines:
            line_stripped = line.strip()
            if not line_stripped:
                # empty line => flush bullets
                flush_bullets()
                continue

            # check for h1
            h1_match = re.match(r'^\*h1-\s+(.*)$', line_stripped, flags=re.IGNORECASE)
            if h1_match:
                flush_bullets()
                heading_text = transform_bold_italic(h1_match.group(1)).upper()
                flowables.append(Paragraph(heading_text, styleH1))
                flowables.append(Spacer(1, 5))
                continue

            # check for h2
            h2_match = re.match(r'^\*h2-\s+(.*)$', line_stripped, flags=re.IGNORECASE)
            if h2_match:
                flush_bullets()
                heading_text = transform_bold_italic(h2_match.group(1)).title()
                flowables.append(Paragraph(heading_text, styleH2))
                flowables.append(Spacer(1, 5))
                continue

            # check bullet line => starts with "* "
            bullet_match = re.match(r'^\*\s+(.*)$', line_stripped)
            if bullet_match:
                bullet_buffer.append(bullet_match.group(1))
                continue

            # else normal line => flush bullets, then paragraph
            flush_bullets()
            ptext = transform_bold_italic(line_stripped)
            flowables.append(Paragraph(ptext, styleNormal))
            flowables.append(Spacer(1, 5))

        # flush leftover bullets
        flush_bullets()

        # Add images (each on a new page)
        if self.image_files:
            flowables.append(PageBreak())

        for ph, img_path in self.image_files.items():
            flowables.append(Paragraph(f"Image Placeholder {ph}", styleH2))
            flowables.append(Spacer(1, 10))
            flowables.append(Image(img_path, width=400, height=300))  # adjust size as needed
            flowables.append(Spacer(1, 20))
            flowables.append(PageBreak())

        doc.build(flowables)
        return out_file


########################################
# Main entry
########################################

if __name__ == "__main__":
    app = QApplication(sys.argv)
    window = MainWindow()
    window.show()
    sys.exit(app.exec_())
