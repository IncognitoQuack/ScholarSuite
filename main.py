import streamlit as st
import os
import requests
import json
from fpdf import FPDF  # fpdf2 library
from pptx import Presentation
from pptx.util import Inches, Pt
import re
import datetime
import io
from pathlib import Path
from typing import Optional, List, Dict

# --- Application Configuration ---
APP_VERSION = "1.1.1"
OPENROUTER_API_URL = "https://openrouter.ai/api/v1/chat/completions"
LLM_MODEL = "meta-llama/llama-3.3-8b-instruct:free"
YOUR_SITE_URL = os.getenv("YOUR_SITE_URL", "http://localhost:8501")
YOUR_SITE_NAME = os.getenv("YOUR_SITE_NAME", "AI Content Creator Deluxe")


# --- Core Functions ---

def generate_prompt(topic: str, doc_type: str, custom_instructions: str) -> str:
    """
    Builds a structured Markdown prompt for the LLM.
    """
    base_prompt = f"""
You are an AI Content Generation Specialist, tasked with producing a high-quality '{doc_type}' document.
The primary topic is: '{topic}'.

Your output MUST be in precise Markdown format, adhering to the following strict structure:
1.  **Main Title:** Start with a single H1 heading (e.g., '# Document Title').
2.  **Section Headings:** Use H2 headings for major sections (e.g., '## Introduction').
3.  **Sub-Sections:** If needed, use H3 headings (e.g., '### Subtopic'). Avoid deeper nesting.
4.  **Paragraphs:** Separate paragraphs with a single blank line.
5.  **Bullet Points:** Use '-' or '*' for lists, each on its own line.

**Content Requirements:**
-   Must include an Introduction, logically structured body sections, and a Conclusion.
-   Use professional, informative tone suited to a '{doc_type}'.

**Formatting Rules (Strict):**
-   **NO PREAMBLE/POSTAMBLE:** Do NOT include phrases like "Here is your document:".
-   **RAW MARKDOWN ONLY:** Entire response must be raw Markdown.
-   **LINE BREAKS:** Use single blank lines around headings and paragraphs.

Example structure:
# {topic}: A Comprehensive {doc_type}

## Introduction
Brief intro of '{topic}' and the '{doc_type}' objective.

## {{Section 1 Title}}
Discuss first major aspect.
- Bullet point A.
- Bullet point B.

### {{Optional Subsection}}
Additional details. (must be added in order to inc the content lenght)

## {{Section 2 Title}}
In-depth exploration.

## Conclusion
Summarize main points and conclude '{doc_type}'.

Now, generate the '{doc_type}' based on these instructions.
"""
    if custom_instructions.strip():
        base_prompt += (
            f"\n**User's Custom Instructions (Override formatting if needed):**\n"
            f"{custom_instructions.strip()}\n"
        )
    return base_prompt.strip()


def call_openrouter_api(prompt: str, api_key: str) -> Optional[str]:
    """
    Calls the OpenRouter.ai API with the given prompt.
    Returns the generated Markdown string, or None on failure.
    """
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json",
        "HTTP-Referer": YOUR_SITE_URL,
        "X-Title": YOUR_SITE_NAME,
    }
    payload = {
        "model": LLM_MODEL,
        "messages": [{"role": "user", "content": prompt}],
        "temperature": 0.65,
        "max_tokens": 4000,
    }
    try:
        response = requests.post(
            OPENROUTER_API_URL,
            headers=headers,
            json=payload,
            timeout=240  # 4-minute timeout
        )
        response.raise_for_status()
        result = response.json()
        choices = result.get("choices", [])
        if choices and "message" in choices[0] and "content" in choices[0]["message"]:
            return choices[0]["message"]["content"].strip()
        else:
            st.error(
                "API Error: Unexpected response structure. Full response:\n"
                f"```json\n{json.dumps(result, indent=2)}\n```"
            )
            return None
    except requests.exceptions.Timeout:
        st.error("API Error: Request timed out (240 seconds). The model might be busy or unavailable.")
        return None
    except requests.exceptions.HTTPError as e:
        detail = f"HTTP {e.response.status_code}"
        try:
            err_json = e.response.json()
            if "error" in err_json:
                msg = (
                    err_json["error"].get("message")
                    if isinstance(err_json["error"], dict)
                    else err_json["error"]
                )
                detail += f" – {msg}"
            else:
                detail += f" – {e.response.text[:300]}"
        except json.JSONDecodeError:
            detail += f" – Non-JSON response: {e.response.text[:300]}"
        st.error(f"API Error: {detail}")
        return None
    except requests.exceptions.RequestException as e:
        st.error(f"API Network Error: {e}")
        return None
    except json.JSONDecodeError as e:
        st.error(
            "API Error: Unable to decode JSON.\n"
            f"Response text: {response.text if 'response' in locals() else 'None'}\n"
            f"Error: {e}"
        )
        return None
    except Exception as e:
        st.error(f"Unexpected Error: {e}")
        return None


def parse_markdown_content(markdown_text: str) -> List[Dict[str, str]]:
    """
    Parses Markdown into a structured list of items:
    - h1, h2, h3, p (paragraph), bullet (list item)
    """
    if not markdown_text:
        return []

    lines = markdown_text.splitlines()
    structured: List[Dict[str, str]] = []
    buffer: List[str] = []

    def flush_paragraph():
        nonlocal buffer
        if buffer:
            paragraph = "\n".join(buffer).strip()
            if paragraph:
                structured.append({"type": "p", "content": paragraph})
            buffer = []

    for line in lines:
        stripped = line.strip()
        if stripped.startswith("# "):
            flush_paragraph()
            structured.append({"type": "h1", "content": stripped[2:].strip()})
        elif stripped.startswith("## "):
            flush_paragraph()
            structured.append({"type": "h2", "content": stripped[3:].strip()})
        elif stripped.startswith("### "):
            flush_paragraph()
            structured.append({"type": "h3", "content": stripped[4:].strip()})
        elif (stripped.startswith("- ") or stripped.startswith("* ")) and len(stripped) > 2:
            flush_paragraph()
            structured.append({"type": "bullet", "content": stripped[2:].strip()})
        elif stripped == "":
            flush_paragraph()
        else:
            buffer.append(line)

    flush_paragraph()
    return structured


class PDF(FPDF):
    """
    Custom PDF class using fpdf2.
    Includes header, footer, and methods to add structured content.
    """
    def __init__(self, *args, **kwargs):
        super().__init__(*args, **kwargs)
        self.doc_title = "Generated Document"
        self.using_dejavu = False
        self.creation_timestamp = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        # Attempt to load DejaVu fonts for wider character support
        try:
            font_dir = Path(__file__).parent
            self.add_font("DejaVu", "", str(font_dir / "DejaVuSansCondensed.ttf"))
            self.add_font("DejaVu", "B", str(font_dir / "DejaVuSansCondensed-Bold.ttf"))
            self.set_font("DejaVu", "", 12)
            self.using_dejavu = True
        except Exception:
            st.sidebar.warning(
                "⚠️ DejaVu fonts not found. PDF will use Arial (limited character support)."
            )
            self.set_font("Arial", "", 12)

    def set_document_metadata(self, title: str, author: str = "AI Content Creator"):
        self.set_title(title)
        self.set_author(author)
        self.set_creator(f"{YOUR_SITE_NAME} (v{APP_VERSION})")
        self.doc_title = title

    def header(self):
        # Skip header on first page (title page)
        if self.page_no() == 1:
            return
        original = (self.font_family, self.font_style, self.font_size_pt)
        fam = "DejaVu" if self.using_dejavu else "Arial"
        self.set_font(fam, "B", 9)
        self.set_y(10)
        self.cell(0, 10, self.doc_title, 0, 0, "C")
        self.ln(12)
        self.set_font(*original)

    def footer(self):
        original = (self.font_family, self.font_style, self.font_size_pt)
        fam = "DejaVu" if self.using_dejavu else "Arial"
        self.set_font(fam, "I", 8)
        self.set_y(-15)
        left_text = f"Generated: {self.creation_timestamp}"
        right_text = f"Page {self.page_no()} / {{nb}}"
        self.cell(0, 10, left_text, 0, 0, "L")
        self.set_x(self.w - self.r_margin - 30)
        self.cell(30, 10, right_text, 0, 0, "R")
        self.set_font(*original)

    def _active_font(self) -> str:
        return "DejaVu" if self.using_dejavu else "Arial"

    def add_title_page(self, title: str, subtitle: str = ""):
        self.add_page()
        fam = self._active_font()
        self.set_font(fam, "B", 28)
        self.ln(self.h / 4)
        self.multi_cell(0, 15, title, 0, "C")
        self.ln(5)
        if subtitle:
            self.set_font(fam, "", 16)
            self.multi_cell(0, 10, subtitle, 0, "C")
        self.ln(10)
        self.set_font(fam, "", 12)

    def chapter_heading(self, text: str, level: int = 1):
        fam = self._active_font()
        if level == 1:
            self.set_font(fam, "B", 18)
            self.ln(10)
        elif level == 2:
            self.set_font(fam, "B", 16)
            self.ln(7)
        else:  # level == 3
            self.set_font(fam, "B", 14)
            self.ln(5)
        self.multi_cell(0, 8, text, 0, "L")
        self.ln(3)
        self.set_font(fam, "", 12)

    def chapter_paragraph(self, text: str):
        fam = self._active_font()
        self.set_font(fam, "", 12)
        self.multi_cell(0, 7, text)
        self.ln(2)

    def chapter_bullet_point(self, text: str):
        fam = self._active_font()
        self.set_font(fam, "", 12)
        bullet_char = "\u2022" if self.using_dejavu else "-"
        self.set_x(self.l_margin + 5)
        self.multi_cell(0, 7, f"{bullet_char} {text}")
        self.set_x(self.l_margin)
        self.ln(1)


def create_pdf(
    structured_content: List[Dict[str, str]],
    doc_topic: str,
    doc_type_str: str
) -> Optional[bytes]:
    """
    Builds a PDF byte stream from structured content.
    """
    if not structured_content:
        st.warning("No content available for PDF generation.")
        return None

    pdf = PDF(orientation="P", unit="mm", format="A4")
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.set_document_metadata(title=doc_topic, author=f"AI for {doc_type_str}")
    pdf.alias_nb_pages()

    # Determine title page content
    main_title = doc_topic
    content_items = structured_content
    if structured_content and structured_content[0]["type"] == "h1":
        main_title = structured_content[0]["content"]
        content_items = structured_content[1:]

    pdf.add_title_page(main_title, subtitle=f"A {doc_type_str} Document")
    if content_items:
        pdf.add_page()

    for item in content_items:
        if pdf.get_y() > (pdf.h - pdf.b_margin - 25):
            pdf.add_page()
        t = item["type"]
        if t == "h1":
            pdf.chapter_heading(item["content"], level=1)
        elif t == "h2":
            pdf.chapter_heading(item["content"], level=2)
        elif t == "h3":
            pdf.chapter_heading(item["content"], level=3)
        elif t == "p":
            pdf.chapter_paragraph(item["content"])
        elif t == "bullet":
            pdf.chapter_bullet_point(item["content"])

    try:
        # output(dest="S") returns a str; encode as Latin-1 to get bytes
        pdf_str = pdf.output(dest="S")
        return pdf_str.encode("latin-1")
    except Exception as e:
        st.error(f"PDF generation error: {e}")
        return None


def create_pptx(
    structured_content: List[Dict[str, str]],
    doc_topic: str,
    doc_type_str: str
) -> Optional[bytes]:
    """
    Builds a PPTX byte stream from structured content.
    """
    if not structured_content:
        st.warning("No content available for PPTX generation.")
        return None

    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]

    main_title = doc_topic
    items = structured_content
    if structured_content and structured_content[0]["type"] == "h1":
        main_title = structured_content[0]["content"]
        items = structured_content[1:]

    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = main_title
    if slide.placeholders and len(slide.placeholders) > 1:
        slide.placeholders[1].text = f"AI-Generated {doc_type_str} | {datetime.date.today():%B %d, %Y}"

    current_body = None

    for item in items:
        t = item["type"]
        if t in ("h1", "h2"):
            slide = prs.slides.add_slide(content_layout)
            slide.shapes.title.text = item["content"]
            current_body = slide.placeholders[1]
            tf = current_body.text_frame
            tf.clear()
            tf.word_wrap = True
            if not tf.paragraphs:
                tf.add_paragraph()
        elif t == "h3":
            if current_body is None:
                slide = prs.slides.add_slide(content_layout)
                slide.shapes.title.text = "Details"
                current_body = slide.placeholders[1]
                tf = current_body.text_frame
                tf.clear()
                tf.word_wrap = True
                if not tf.paragraphs:
                    tf.add_paragraph()
            p = current_body.text_frame.add_paragraph()
            p.text = item["content"]
            p.font.bold = True
            p.font.size = Pt(20)
            p.level = 1
        elif t == "p":
            if current_body is None:
                slide = prs.slides.add_slide(content_layout)
                slide.shapes.title.text = "Content"
                current_body = slide.placeholders[1]
                tf = current_body.text_frame
                tf.clear()
                tf.word_wrap = True
                if not tf.paragraphs:
                    tf.add_paragraph()
            p = current_body.text_frame.add_paragraph()
            p.text = item["content"]
            p.font.size = Pt(18)
            p.level = 0
        elif t == "bullet":
            if current_body is None:
                slide = prs.slides.add_slide(content_layout)
                slide.shapes.title.text = "Key Points"
                current_body = slide.placeholders[1]
                tf = current_body.text_frame
                tf.clear()
                tf.word_wrap = True
                if not tf.paragraphs:
                    tf.add_paragraph()
            p = current_body.text_frame.add_paragraph()
            p.text = item["content"]
            p.font.size = Pt(18)
            p.level = 1

    try:
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()
    except Exception as e:
        st.error(f"PPTX generation error: {e}")
        return None


# --- Streamlit UI ---

st.set_page_config(
    page_title="AI Content Creator Deluxe",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Initialize session state variables
state = st.session_state
state.setdefault("generated_markdown", None)
state.setdefault("pdf_bytes", None)
state.setdefault("pptx_bytes", None)
state.setdefault("file_basename", "ai_document")
state.setdefault("error_message", None)
state.setdefault("generation_complete", False)

# Sidebar
with st.sidebar:
    st.image("https://cdn-icons-png.flaticon.com/512/2593/2593194.png", width=69)
    st.title(YOUR_SITE_NAME)
    st.caption(f"Version: {APP_VERSION}")
    st.markdown("---")

    st.subheader("🔑 API Configuration")
    api_input = st.text_input(
        "OpenRouter.ai API Key:",
        type="password",
        placeholder="sk-or-...",
        help="Get your API key from https://openrouter.ai/keys"
    )
    OPENROUTER_API_KEY = api_input.strip() or os.getenv("OPENROUTER_API_KEY", "")
    if OPENROUTER_API_KEY:
        st.success("API Key loaded.")
        st.caption(f"Using model: `{LLM_MODEL}`")
    else:
        st.warning("API Key is required. Enter it above or set `OPENROUTER_API_KEY` env variable.")

    st.markdown("---")
    st.subheader("📝 How to Use")
    st.markdown(
        """
1. Enter your OpenRouter.ai API Key.
2. Define **Topic** and select **Document Type**.
3. (Optional) Add **Custom Instructions** to guide AI.
4. Click **Generate Content**.
5. Download your PDF / PPTX when ready.
"""
    )
    # st.caption("For best PDF fonts, place DejaVu `.ttf` files next to this script.")

# Main area
st.header("🤖 AI Content Creator Deluxe ✨")
st.markdown("Transform your ideas into structured documents (PDFs & PPTX) using AI.")
st.markdown("---")

col_form, col_output = st.columns([0.45, 0.55], gap="large")

with col_form:
    st.subheader("📝 Document Specifications")
    with st.form("generation_form", clear_on_submit=False):
        topic = st.text_input(
            "Topic / Main Subject:",
            placeholder="e.g., The Future of Quantum Computing",
            help="Central theme of your document."
        )
        doc_types = [
            "Article", "Blog Post", "Report", "Case Study", "Project Description",
            "Product Description", "Essay", "Speech Outline", "Research Summary",
            "Technical Documentation", "Marketing Copy", "Meeting Minutes Outline"
        ]
        doc_type = st.selectbox(
            "Document Type:",
            doc_types,
            index=0,
            help="Select the document format."
        )
        custom_instructions = st.text_area(
            "Custom Instructions (Optional):",
            placeholder="e.g., Tone: formal. Include 3 statistics. Max 500-word introduction.",
            height=150,
            help="Fine-tune AI output (tone, sections, etc.)."
        )
        submit = st.form_submit_button("🚀 Generate Content", use_container_width=True)

    if submit:
        # Reset state
        state.generated_markdown = None
        state.pdf_bytes = None
        state.pptx_bytes = None
        state.error_message = None
        state.generation_complete = False

        if not OPENROUTER_API_KEY:
            st.error("🔴 Error: OpenRouter.ai API Key is missing.")
        elif not topic.strip():
            st.error("🔴 Error: Topic cannot be empty.")
        else:
            # Sanitize filename
            sanitized = re.sub(r"[^\w\s-]", "", topic.lower()).strip()
            sanitized = re.sub(r"[-\s]+", "_", sanitized)
            timestamp = datetime.datetime.now().strftime("%Y%m%d%H%M%S")
            state.file_basename = f"{doc_type.replace(' ', '_')}_{sanitized[:40]}_{timestamp}"

            progress = st.progress(0, text="Initializing...")
            status = st.empty()

            try:
                status.info("⚙️ Phase 1/5: Building prompt...")
                prompt = generate_prompt(topic.strip(), doc_type, custom_instructions)
                progress.progress(20, text="Prompt ready. Contacting AI...")

                status.info("🧠 Phase 2/5: Querying AI model...")
                with st.spinner("Waiting for AI response..."):
                    raw_md = call_openrouter_api(prompt, OPENROUTER_API_KEY)
                progress.progress(50, text="AI response received.")

                if not raw_md:
                    state.error_message = "AI returned no content."
                else:
                    state.generated_markdown = raw_md
                    status.success("✅ AI content received. Parsing...")
                    structured = parse_markdown_content(raw_md)
                    progress.progress(65, text="Content parsed.")

                    if not structured:
                        state.error_message = "Parsed content is empty or invalid."
                        st.warning("AI output could not be parsed into structure.")
                    else:
                        # Generate PDF
                        status.info("📄 Phase 3/5: Generating PDF...")
                        with st.spinner("Creating PDF..."):
                            pdf_out = create_pdf(
                                structured,
                                doc_topic=topic.strip(),
                                doc_type_str=doc_type
                            )
                        if pdf_out:
                            state.pdf_bytes = pdf_out
                            progress.progress(80, text="PDF ready.")
                        else:
                            state.error_message = (
                                (state.error_message or "") + " | PDF generation failed."
                            )

                        # Generate PPTX
                        status.info("🖥️ Phase 4/5: Generating PPTX...")
                        with st.spinner("Creating PPTX..."):
                            pptx_out = create_pptx(
                                structured,
                                doc_topic=topic.strip(),
                                doc_type_str=doc_type
                            )
                        if pptx_out:
                            state.pptx_bytes = pptx_out
                            progress.progress(95, text="PPTX ready.")
                        else:
                            state.error_message = (
                                (state.error_message or "") + " | PPTX generation failed."
                            )

                status.info("🏁 Phase 5/5: Finalizing...")
                state.generation_complete = True
                progress.progress(100, text="Done.")
                if state.pdf_bytes or state.pptx_bytes:
                    status.success("🎉 Documents are ready for download below.")
                else:
                    status.warning("Process finished but no documents were created.")

            except Exception as e:
                st.error(f"🆘 Critical error: {e}")
                st.exception(e)
                state.error_message = f"Critical error: {e}"
                state.generation_complete = True
                progress.progress(100)
                status.error("Process halted due to an unexpected error.")

with col_output:
    st.subheader("📂 Generated Outputs & Preview")

    if state.generation_complete:
        if state.error_message and not (state.pdf_bytes or state.pptx_bytes):
            st.error(f"🔴 Generation failed. Details: {state.error_message}")

        if state.pdf_bytes:
            st.success("✅ PDF Document Generated!")
            st.download_button(
                "📥 Download PDF",
                data=state.pdf_bytes,
                file_name=f"{state.file_basename}.pdf",
                mime="application/pdf",
                use_container_width=True,
                key="download_pdf"
            )
            st.markdown("---")

        if state.pptx_bytes:
            st.success("✅ PowerPoint Presentation Generated!")
            st.download_button(
                "📥 Download PPTX",
                data=state.pptx_bytes,
                file_name=f"{state.file_basename}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
                key="download_pptx"
            )
            st.markdown("---")

        if state.generated_markdown:
            with st.expander("👁️ View Raw Markdown Content", expanded=False):
                st.code(state.generated_markdown, language="markdown")
        elif state.error_message:
            st.info("No Markdown preview due to errors.")
        else:
            st.info("Generation complete, but no content to preview.")
    else:
        st.info("After you click 'Generate Content', your documents and preview will appear here.")

st.markdown("---")
st.caption(f"AI Content Creator Deluxe v{APP_VERSION} | Beta.")
